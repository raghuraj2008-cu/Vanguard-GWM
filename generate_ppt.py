from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Initialize Presentation & 16:9 Widescreen Dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
C_BG = RGBColor(245, 247, 250)         # Light Grey Background
C_CARD = RGBColor(255, 255, 255)       # Pure White Card Fill
C_TEXT = RGBColor(30, 41, 59)          # Dark Slate Text
C_HEADER_BLUE = RGBColor(15, 37, 55)   # Deep Navy Header
C_ACCENT_BLUE = RGBColor(14, 116, 144) # Professional Cyan/Blue Accent
C_SHAPE_BG = RGBColor(230, 238, 248)   # Light Accent Shape Fill

blank_slide_layout = prs.slide_layouts[6]

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG

def add_sih_header(slide, title_text, slide_num):
    # 1. Top-Left: APEX CODERS Header Badge Shape
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.35), Inches(2.2), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_SHAPE_BG
    badge.line.color.rgb = C_ACCENT_BLUE
    badge.line.width = Pt(1.5)
    
    tb_badge = slide.shapes.add_textbox(Inches(0.8), Inches(0.37), Inches(2.2), Inches(0.5))
    tf_b = tb_badge.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = "APEX CODERS"
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = C_HEADER_BLUE
    p_b.alignment = PP_ALIGN.CENTER

    # 2. Center: Slide Title Shape Banner
    title_box_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), Inches(0.35), Inches(6.4), Inches(0.55))
    title_box_shape.fill.solid()
    title_box_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title_box_shape.line.color.rgb = RGBColor(203, 213, 225)
    title_box_shape.line.width = Pt(1)

    tb_title = slide.shapes.add_textbox(Inches(3.3), Inches(0.38), Inches(6.2), Inches(0.5))
    tf_t = tb_title.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title_text.upper()
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = C_HEADER_BLUE
    p_t.alignment = PP_ALIGN.CENTER

    # 3. Top-Right: SIH Branding Box Shape
    logo_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.35), Inches(2.733), Inches(0.55))
    logo_shape.fill.solid()
    logo_shape.fill.fore_color.rgb = RGBColor(230, 238, 248)
    logo_shape.line.color.rgb = C_ACCENT_BLUE
    logo_shape.line.width = Pt(1.5)

    tb_logo = slide.shapes.add_textbox(Inches(9.8), Inches(0.38), Inches(2.733), Inches(0.5))
    tf_l = tb_logo.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "SMART INDIA HACKATHON 2026"
    p_l.font.size = Pt(9)
    p_l.font.bold = True
    p_l.font.color.rgb = C_HEADER_BLUE
    p_l.alignment = PP_ALIGN.CENTER

    # 4. Bottom Footer Banner
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(255, 255, 255)
    footer_bar.line.color.rgb = RGBColor(203, 213, 225)

    tb_foot = slide.shapes.add_textbox(Inches(1.0), Inches(7.0), Inches(11.3), Inches(0.35))
    tf_f = tb_foot.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = f"@SIH Idea submission - Template                                                                                                                                              {slide_num}"
    p_f.font.size = Pt(9)
    p_f.font.color.rgb = RGBColor(120, 130, 145)

def add_content_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD
    shape.line.color.rgb = RGBColor(203, 213, 225)
    shape.line.width = Pt(1.5)
    return shape

# ==========================================
# SLIDE 1: TITLE PAGE
# ==========================================
s1 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s1)
add_sih_header(s1, "Title Page", 1)
add_content_card(s1, 0.8, 1.1, 11.733, 5.7)

# Left Card: Team Members
box_tm = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.3), Inches(5.5), Inches(5.3))
box_tm.fill.solid()
box_tm.fill.fore_color.rgb = RGBColor(255, 255, 255)
box_tm.line.color.rgb = C_ACCENT_BLUE
box_tm.line.width = Pt(1)

tb_tm = s1.shapes.add_textbox(Inches(1.25), Inches(1.45), Inches(5.0), Inches(5.0))
tf_tm = tb_tm.text_frame
tf_tm.word_wrap = True
p_tmh = tf_tm.paragraphs[0]
p_tmh.text = "TEAM MEMBERS"
p_tmh.font.size = Pt(12)
p_tmh.font.bold = True
p_tmh.font.color.rgb = C_HEADER_BLUE
p_tmh.space_after = Pt(10)

members = [
    ("Raghuraj Pratap Singh", "Team Leader"),
    ("Ayush Singh", "Member"),
    ("Himanshu Kumar", "Member"),
    ("Rohit Gour", "Member"),
    ("Astha Soni", "Member"),
    ("Simran Jaiswal", "Member")
]
for name, role in members:
    p = tf_tm.add_paragraph()
    p.text = f"{name} ({role})"
    p.font.size = Pt(10)
    p.font.color.rgb = C_TEXT
    p.space_after = Pt(6)

# Right Card: Project Metadata
box_ps = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.3), Inches(5.4), Inches(5.3))
box_ps.fill.solid()
box_ps.fill.fore_color.rgb = RGBColor(250, 250, 250)
box_ps.line.color.rgb = RGBColor(203, 213, 225)
box_ps.line.width = Pt(1)

tb_ps = s1.shapes.add_textbox(Inches(7.05), Inches(1.45), Inches(5.1), Inches(5.0))
tf_ps = tb_ps.text_frame
tf_ps.word_wrap = True
p_psh = tf_ps.paragraphs[0]
p_psh.text = "PROJECT METADATA"
p_psh.font.size = Pt(12)
p_psh.font.bold = True
p_psh.font.color.rgb = C_HEADER_BLUE
p_psh.space_after = Pt(10)

meta_items = [
    ("Problem Statement ID:", "SIH26-PS153"),
    ("Problem Statement Title:", "Graph World Model for Predictive Cyber Attack Forecasting"),
    ("Theme:", "Cybersecurity & Critical Information Infrastructure"),
    ("PS Category:", "Software"),
    ("Team ID & Name:", "SIH2026 PS153 APEX / APEX CODERS")
]
for mk, mv in meta_items:
    p = tf_ps.add_paragraph()
    p.text = f"• {mk} {mv}"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = C_TEXT
    p.space_after = Pt(6)

# ==========================================
# SLIDE 2: IDEA TITLE & PROPOSED SOLUTION
# ==========================================
s2 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s2)
add_sih_header(s2, "Idea Title: Vanguard-GWM Platform", 2)
add_content_card(s2, 0.8, 1.1, 11.733, 5.7)

tb2 = s2.shapes.add_textbox(Inches(1.1), Inches(1.3), Inches(11.1), Inches(5.3))
tf2 = tb2.text_frame
tf2.word_wrap = True

s2_points = [
    ("Proposed Solution: Vanguard-GWM", "Models computer networks as evolving temporal graphs to forecast multi-stage APT attacks before host compromise occurs."),
    ("Core Pipeline & Workflow", "• Ingests dual-level telemetry (flow aggregates + transport header geometry). • Constructs dynamic network graph G_t using Polars and NetworkX. • Employs PyTorch Geometric GATv2 spatial encoders and Causal Temporal Transformers to predict state transitions and execute K-step autoregressive rollouts."),
    ("How it Addresses the Core Challenge", "• Replaces static, reactive classifiers (Random Forest / LSTM) that trigger post-facto. • Delivers +180s early warning during the reconnaissance phase, stopping multi-stage APTs prior to lateral movement."),
    ("Innovation & Uniqueness", "• True predictive world model rather than a reactive signature filter. • Payload-agnostic telemetry fusion with built-in Explainability (XAI) and What-If containment simulation.")
]
for h, b in s2_points:
    p = tf2.add_paragraph()
    p.text = f"• {h}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT_BLUE
    
    p2 = tf2.add_paragraph()
    p2.text = f"  {b}"
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = C_TEXT
    p2.space_after = Pt(4)

# ==========================================
# SLIDE 3: TECHNICAL APPROACH
# ==========================================
s3 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s3)
add_sih_header(s3, "Technical Approach", 3)
add_content_card(s3, 0.8, 1.1, 11.733, 5.7)

steps_data = [
    ("1. Telemetry Ingestion", "Scapy, PyShark,\nPolars (Δt=1-5s)"),
    ("2. Graph Construction", "Dynamic Graph G_t\n(Nodes & Edges)"),
    ("3. World Model Rollout", "GATv2 + Causal\nTransformer (K-steps)"),
    ("4. SDN Actuation", "OpenFlow 1.3\nQUARANTINE / DROP")
]
box_width = 2.6
box_left_start = 1.1
for i, (st_title, st_desc) in enumerate(steps_data):
    bx = box_left_start + (i * 2.85)
    box_shape = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(1.3), Inches(box_width), Inches(1.1))
    box_shape.fill.solid()
    box_shape.fill.fore_color.rgb = C_SHAPE_BG
    box_shape.line.color.rgb = C_ACCENT_BLUE
    box_shape.line.width = Pt(1)
    
    tb_box = s3.shapes.add_textbox(Inches(bx), Inches(1.35), Inches(box_width), Inches(1.0))
    tf_bx = tb_box.text_frame
    tf_bx.word_wrap = True
    p_bt = tf_bx.paragraphs[0]
    p_bt.text = st_title
    p_bt.font.size = Pt(9)
    p_bt.font.bold = True
    p_bt.font.color.rgb = C_HEADER_BLUE
    p_bt.alignment = PP_ALIGN.CENTER
    
    p_bd = tf_bx.add_paragraph()
    p_bd.text = st_desc
    p_bd.font.size = Pt(8)
    p_bd.font.color.rgb = C_TEXT
    p_bd.alignment = PP_ALIGN.CENTER

tb3 = s3.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(11.1), Inches(4.1))
tf3 = tb3.text_frame
tf3.word_wrap = True

s3_points = [
    ("Methodology & Implementation Flow:", "• Continuous streaming telemetry via Scapy/PyShark/Polars. Extracts rolling window metrics to construct dynamic graph G_t. • PyTorch Geometric GATv2 encoder compresses network topology into latent state z_t; Causal Temporal Transformer learns transition dynamics P(z_{t+1}|z_t) [World Model core]. • Performs K-step forward simulation predicting risk score and staging automated OpenFlow 1.3 SDN DROP rules via a Streamlit SOC interface."),
    ("Technology Stack", "• Frontend/UI: Streamlit, Plotly | Backend: Python FastAPI | AI/ML Framework: PyTorch, PyTorch Geometric | Database: NetworkX, Polars Dataframes | Networking API: OpenFlow 1.3")
]
for h, b in s3_points:
    p = tf3.add_paragraph()
    p.text = f"• {h}"
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT_BLUE
    if b.strip():
        p2 = tf3.add_paragraph()
        p2.text = f"  {b}"
        p2.font.size = Pt(9)
        p2.font.color.rgb = C_TEXT
    p.space_after = Pt(4)

# ==========================================
# SLIDE 4: FEASIBILITY AND VIABILITY
# ==========================================
s4 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s4)
add_sih_header(s4, "Feasibility and Viability", 4)
add_content_card(s4, 0.8, 1.1, 11.733, 5.7)

risk_shape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.3), Inches(5.5), Inches(5.3))
risk_shape.fill.solid()
risk_shape.fill.fore_color.rgb = RGBColor(255, 242, 242)
risk_shape.line.color.rgb = RGBColor(239, 68, 68)
risk_shape.line.width = Pt(1.2)

tb_r = s4.shapes.add_textbox(Inches(1.25), Inches(1.45), Inches(5.2), Inches(5.0))
tf_r = tb_r.text_frame
tf_r.word_wrap = True
p_rh = tf_r.paragraphs[0]
p_rh.text = "POTENTIAL CHALLENGES & RISKS"
p_rh.font.size = Pt(11)
p_rh.font.bold = True
p_rh.font.color.rgb = RGBColor(185, 28, 28)
p_rh.space_after = Pt(6)

risks_list = [
    ("Telemetry Parsing Latency:", "Slow PCAP processing during live demo. Fix: Pre-extract feature matrices ahead of time and use Polars for high-speed streaming dataframes."),
    ("Graph Over-smoothing:", "GAT layers losing separation on large networks. Fix: Use a 2-3 layer GAT with InfoNCE contrastive loss over subgraphs."),
    ("Autoregressive Latent Drift:", "Drift accumulating beyond K>5 steps. Fix: Cap rollout horizon at K=3-5 time windows and add temporal consistency regularizers.")
]
for rh, rb in risks_list:
    p = tf_r.add_paragraph()
    p.text = f"• {rh}"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = C_TEXT
    p2 = tf_r.add_paragraph()
    p2.text = f"  {rb}"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = RGBColor(70, 80, 95)
    p2.space_after = Pt(3)

feas_shape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.3), Inches(5.4), Inches(5.3))
feas_shape.fill.solid()
feas_shape.fill.fore_color.rgb = RGBColor(240, 253, 244)
feas_shape.line.color.rgb = RGBColor(16, 185, 129)
feas_shape.line.width = Pt(1.2)

tb_f = s4.shapes.add_textbox(Inches(7.05), Inches(1.45), Inches(5.1), Inches(5.0))
tf_f = tb_f.text_frame
tf_f.word_wrap = True
p_fh = tf_f.paragraphs[0]
p_fh.text = "FEASIBILITY & PROTOTYPE SCOPE"
p_fh.font.size = Pt(11)
p_fh.font.bold = True
p_fh.font.color.rgb = RGBColor(4, 120, 87)
p_fh.space_after = Pt(6)

feas_list = [
    ("Verified Performance:", "Fully prototyped with sub-50ms inference latency (36.4 ms), operating entirely offline/air-gapped without cloud API dependencies."),
    ("Prototype Scope Pipeline:", "Telemetry Ingestion ➔ OCSF/Feature Slicing ➔ GATv2 Encoder ➔ World Model Core ➔ SDN Enforcer."),
    ("Deployment Readiness:", "Pinned Docker environment, automated unit testing suite, and complete GitHub repository structure.")
]
for fh, fb in feas_list:
    p = tf_f.add_paragraph()
    p.text = f"• {fh}"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = C_TEXT
    p2 = tf_f.add_paragraph()
    p2.text = f"  {fb}"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = RGBColor(70, 80, 95)
    p2.space_after = Pt(4)

# ==========================================
# SLIDE 5: IMPACT AND BENEFITS
# ==========================================
s5 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s5)
add_sih_header(s5, "Impact and Benefits", 5)
add_content_card(s5, 0.8, 1.1, 11.733, 5.7)

box_imp_l = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.3), Inches(5.5), Inches(5.3))
box_imp_l.fill.solid()
box_imp_l.fill.fore_color.rgb = RGBColor(238, 242, 255)
box_imp_l.line.color.rgb = RGBColor(99, 102, 241)
box_imp_l.line.width = Pt(1.2)

tb_il = s5.shapes.add_textbox(Inches(1.25), Inches(1.45), Inches(5.2), Inches(5.0))
tf_il = tb_il.text_frame
tf_il.word_wrap = True
p_ilh = tf_il.paragraphs[0]
p_ilh.text = "POTENTIAL IMPACT ON TARGET AUDIENCE"
p_ilh.font.size = Pt(11)
p_ilh.font.bold = True
p_ilh.font.color.rgb = RGBColor(79, 70, 229)
p_ilh.space_after = Pt(6)

impact_items = [
    ("Enterprise SOCs & CII:", "Shifts cybersecurity operations completely from reactive alert triage to proactive, predictive defense."),
    ("Eliminates Alert Fatigue:", "Groups thousands of fragmented log entries into a single, cohesive kill-chain trajectory narrative."),
    ("National & Defense Impact:", "Secures critical national infrastructure against stealthy, zero-day Advanced Persistent Threats (APTs).")
]
for ih, ib in impact_items:
    p = tf_il.add_paragraph()
    p.text = f"• {ih}"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = C_TEXT
    p2 = tf_il.add_paragraph()
    p2.text = f"  {ib}"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = RGBColor(70, 80, 95)
    p2.space_after = Pt(4)

box_imp_r = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.3), Inches(5.4), Inches(5.3))
box_imp_r.fill.solid()
box_imp_r.fill.fore_color.rgb = RGBColor(240, 253, 244)
box_imp_r.line.color.rgb = RGBColor(16, 185, 129)
box_imp_r.line.width = Pt(1.2)

tb_ir = s5.shapes.add_textbox(Inches(7.05), Inches(1.45), Inches(5.1), Inches(5.0))
tf_ir = tb_ir.text_frame
tf_ir.word_wrap = True
p_irh = tf_ir.paragraphs[0]
p_irh.text = "BENEFITS & QUANTIFIED OUTCOMES"
p_irh.font.size = Pt(11)
p_irh.font.bold = True
p_irh.font.color.rgb = RGBColor(4, 120, 87)
p_irh.space_after = Pt(6)

benefit_items = [
    ("Proactive Lead Time:", "+180s to +600s early warning during initial reconnaissance with 97.9% multi-stage APT recall[cite: 1] (2.0% FPR[cite: 1])."),
    ("Decision Support:", "What-If simulation lets analysts test host quarantine actions (ΔRisk = -61.1%[cite: 1]) before acting."),
    ("Economic Value:", "Protects critical infrastructure, minimizes breach costs, and prevents analyst burnout.")
]
for bh, bb in benefit_items:
    p = tf_ir.add_paragraph()
    p.text = f"• {bh}"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = C_TEXT
    p2 = tf_ir.add_paragraph()
    p2.text = f"  {bb}"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = RGBColor(70, 80, 95)
    p2.space_after = Pt(4)

# ==========================================
# SLIDE 6: RESEARCH AND REFERENCES (4-Box Grid)
# ==========================================
s6 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(s6)
add_sih_header(s6, "Research and References", 6)
add_content_card(s6, 0.8, 1.1, 11.733, 5.7)

grid_boxes = [
    (1.1, 1.3, "PRIMARY BENCHMARK DATASETS", "• CSE-CIC-IDS2018 and CTU-13 network traffic datasets.\n• Attack timelines divided into temporal windows to build sequential network states and ground-truth state transitions."),
    (6.9, 1.3, "EVALUATION METHODOLOGY", "• Held-out attack variants to measure unseen-attack generalization.\n• Benchmarked against Logistic Regression & Random Forest baselines across Precision, Recall, F1, and Predictive Lead Time[cite: 1]."),
    (1.1, 4.1, "CORE ACADEMIC REFERENCES", "• Ha, D., & Schmidhuber, J. (2018). 'World Models.' arXiv:1803.10122.\n• Velickovic, P., et al. (2018). 'Graph Attention Networks.' ICLR."),
    (6.9, 4.1, "PROJECT REPOSITORY & DOCS", "• GitHub Repository: github.com/raghuraj2008-cu/Vanguard-GWM[cite: 1]\n• IEEE Conference Manuscript: Vanguard_GWM_IEEE_Conference_Paper.pdf[cite: 1]")
]

for gx, gy, gtitle, gdesc in grid_boxes:
    g_shape = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(gx), Inches(gy), Inches(5.4), Inches(2.6))
    g_shape.fill.solid()
    g_shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
    g_shape.line.color.rgb = RGBColor(203, 213, 225)
    g_shape.line.width = Pt(1)
    
    tb_g = s6.shapes.add_textbox(Inches(gx + 0.15), Inches(gy + 0.1), Inches(5.1), Inches(2.4))
    tf_g = tb_g.text_frame
    tf_g.word_wrap = True
    
    pg_t = tf_g.paragraphs[0]
    pg_t.text = gtitle
    pg_t.font.size = Pt(10)
    pg_t.font.bold = True
    pg_t.font.color.rgb = C_HEADER_BLUE
    pg_t.space_after = Pt(3)
    
    pg_d = tf_g.add_paragraph()
    pg_d.text = gdesc
    pg_d.font.size = Pt(8.5)
    pg_d.font.color.rgb = C_TEXT

# Save Presentation
prs.save("CYVERGE_SIH_2026_Vanguard_GWM.pptx")
print("✅ CYVERGE_SIH_2026_Vanguard_GWM.pptx generated successfully with fully visible text content!")
